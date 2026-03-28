#!/usr/bin/env python3
"""Build BestPath hackathon presentation as PPTX with speaker notes and images."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "images")

# Teal brand colours
TEAL = RGBColor(0x0B, 0x85, 0x85)
NAVY = RGBColor(0x1A, 0x2E, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM_WHITE = RGBColor(0xFA, 0xFA, 0xF8)
LIGHT_TEAL = RGBColor(0xE0, 0xF2, 0xF1)
DARK_BG = RGBColor(0x12, 0x20, 0x2E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def img_path(name):
    return os.path.join(IMG, name)


def add_image_slide(prs, image_file, notes_text, title_text=None, subtitle_text=None):
    """Add a slide with a full-bleed image and speaker notes."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Set slide background to dark
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Calculate image dimensions to fill slide (cover)
    ipath = img_path(image_file)
    with Image.open(ipath) as im:
        iw, ih = im.size

    slide_ratio = 13.333 / 7.5
    img_ratio = iw / ih

    if img_ratio >= slide_ratio:
        # Image is wider — fit height, crop width
        h = SLIDE_H
        w = Emu(int(SLIDE_H * img_ratio))
        left = Emu(int((SLIDE_W - w) / 2))
        top = Emu(0)
    else:
        # Image is taller — fit width, crop height
        w = SLIDE_W
        h = Emu(int(SLIDE_W / img_ratio))
        left = Emu(0)
        top = Emu(int((SLIDE_H - h) / 2))

    slide.shapes.add_picture(ipath, left, top, w, h)

    # Optional overlay title
    if title_text:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p2.alignment = PP_ALIGN.LEFT

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

    return slide


def add_text_slide(prs, title, body_lines, notes_text, accent_color=TEAL):
    """Add a styled text slide with bullet points or paragraphs."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Warm white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WARM_WHITE

    # Left accent bar
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Emu(0), Emu(0), Inches(0.15), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Body
    body_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.8))
    btf = body_box.text_frame
    btf.word_wrap = True

    for i, line in enumerate(body_lines):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()

        # Handle bold headers like "**Header:** rest"
        if line.startswith("**") and ":**" in line:
            bold_part = line.split(":**")[0].replace("**", "")
            rest = line.split(":**")[1].strip()
            run_b = p.add_run()
            run_b.text = bold_part + ": "
            run_b.font.size = Pt(20)
            run_b.font.bold = True
            run_b.font.color.rgb = TEAL
            run_r = p.add_run()
            run_r.text = rest
            run_r.font.size = Pt(20)
            run_r.font.color.rgb = NAVY
        elif line.startswith("- "):
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = NAVY
            p.level = 1
        elif line == "":
            p.text = ""
            p.font.size = Pt(10)
        else:
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = NAVY

        p.space_after = Pt(6)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

    return slide


def add_bold_text_slide(prs, lines_with_formatting, notes_text, bg_color=WARM_WHITE):
    """Add a text slide with mixed bold/normal formatting per line."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Left accent bar
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Inches(0.15), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.3), Inches(5.5))
    btf = body_box.text_frame
    btf.word_wrap = True

    for i, (text, is_bold, font_size, color) in enumerate(lines_with_formatting):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = is_bold
        run.font.color.rgb = color
        p.space_after = Pt(8)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ─────────────────────────────────────────────
    # SLIDE 1 — Opening selfie (calm room)
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image17.jpg",
        notes_text="Hi everyone.")

    # ─────────────────────────────────────────────
    # SLIDE 2 — Party selfie
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image4.png",
        notes_text=(
            "This has been really fun -- has everyone else had a good time?\n\n"
            "It's been great meeting some of you, and going hard on a project together.\n\n"
            "Who was still working after midnight?\n\n"
            "After 1am?\n\n"
            "1am is about when we sent our detailed spec to Claude to work overnight.\n\n"
            "So, who are we?"
        ))

    # ─────────────────────────────────────────────
    # SLIDE 3 — Childhood brothers photo
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image20.jpg",
        notes_text="We are brothers -- and also a 10-month-old living room tech startup.")

    # ─────────────────────────────────────────────
    # SLIDE 4 — Living room "office"
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image8.jpg",
        notes_text="This is the living room. And the office. Same room.")

    # ─────────────────────────────────────────────
    # SLIDE 5 — PC build / server
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image1.jpg",
        notes_text="And this is our server. Hand-built with love and a concerning number of zip ties.")

    # ─────────────────────────────────────────────
    # SLIDE 6 — Capital Reasoning booth
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image13.jpg",
        notes_text=(
            "We quickly want to thank Lautaro for making this happen.\n\n"
            "Thank you to the judges.\n\n"
            "Thank you to BuildersVault -- very cool.\n\n"
            "And thank you to our sponsors: Red Bull, and Trelent -- Calum.\n\n"
            "Okay, here we go. Hackathon."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 7 — Deadlines: juggling everything
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image18.png",
        notes_text=(
            "When you're faced with a deadline -- whether that's for moving out, "
            "planning a wedding, or doing the dishes -- large or small, tasks can only "
            "be done one box, one detail, or one dish at a time.\n\n"
            "And if you don't take the next step, it can lead to genuine disaster."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 8 — Buried / overwhelmed
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image2.png",
        notes_text=(
            "The person moving in finds you and all your stuff. "
            "The guests at your wedding don't have cutlery. "
            "Or a grumpy housemate has to do the dishes themselves."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 9 — RIP gravestone (health consequences)
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image12.png",
        notes_text=(
            "When it comes to your health, not taking the next step could mean "
            "ending up in a serious medical emergency -- for something that would "
            "have been entirely preventable if acted upon six months earlier."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 10 — Health data meditation
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image7.png",
        notes_text=(
            "Here's the thing: it's not a problem of quantity or quality of health "
            "information. It's a problem of utilization.\n\n"
            "The information is out there -- authoritative, publicly available medical "
            "literature that outlines risk factors, screening intervals, and preventative actions.\n\n"
            "Patient information is also out there -- well protected, but under-utilized "
            "for proactive measures.\n\n"
            "Even people without a primary care provider know their basic metrics, "
            "history, and medications."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 11 — Public Evidence + Patient Data orbs
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image14.png",
        notes_text=(
            "One of the most profound capabilities of AI is reliably connecting "
            "and generating insight from vast amounts of information that, before, "
            "was completely impractical to synthesize.\n\n"
            "So how do we leverage that to address the issues in the healthcare system?"
        ))

    # ─────────────────────────────────────────────
    # SLIDE 12 — Healthcare Venn diagram
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image3.png",
        notes_text=(
            "We know that primary care providers are overloaded, that non-emergency "
            "health needs get routed into high-cost emergency pathways, and that many "
            "patients are without primary care at all.\n\n"
            "Where these fractures overlap, it compounds into a complex, systemic "
            "issue that's hard to overcome."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 13 — Venn + stressed doctor
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image21.png",
        notes_text=(
            "So if those are the issues we're working to overcome, and we have "
            "this amazing technology for generating insight -- how do we do that?"
        ))

    # ─────────────────────────────────────────────
    # SLIDE 14 — BestPath pill bottle reveal
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image22.png",
        notes_text="This is where BestPath comes in.")

    # ─────────────────────────────────────────────
    # SLIDE 15 — What is BestPath (text)
    # ─────────────────────────────────────────────
    add_text_slide(prs,
        title="BestPath",
        body_lines=[
            "A proactive care engine for patients, providers,",
            "people, and the health system as a whole.",
            "",
            "It compares patient data points with clinical guidelines",
            "to surface their greatest risks, and provides one clear",
            "next action based on authoritative medical guides:",
            "",
            "**What:** should be done",
            "**When:** it should be done",
            "**Where:** care can happen",
        ],
        notes_text=(
            "BestPath is a proactive care engine -- for patients, for providers, "
            "for people, and for the health system as a whole.\n\n"
            "It compares patient data points with clinical guidelines to surface "
            "their greatest risks, and provides one clear next action based on "
            "existing authoritative medical guides: what should be done, when it "
            "should be done, and where care can happen."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 16 — One click operational (text)
    # ─────────────────────────────────────────────
    add_text_slide(prs,
        title="One Click. Operational.",
        body_lines=[
            "**For the clinician:** a one-page review and a click.",
            "",
            "**For the patient:** a text and an email.",
            "",
            "It finds the single highest-value preventative or proactive",
            "clinical action for each patient -- and for people who need",
            "to take action, it's made operational with one click.",
        ],
        notes_text=(
            "For the clinician, it's a one-page review and a click. "
            "For the patient, it's a text and an email.\n\n"
            "It finds the single highest-value preventative or proactive "
            "clinical action for each patient, and for people who need to "
            "take action, it's made operational with one click."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 17 — RCV / SRV Engine infographic
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image11.png",
        notes_text=(
            "Risk Confluence Value -- who is this patient? "
            "It combines demographics, active conditions, vitals, medications, "
            "and utilization patterns into one risk profile. Derives key risk states. "
            "Assigns confidence based on data quality and recency.\n\n"
            "Screening Recency Value -- what has been done, and when? "
            "Tracks each relevant action's last completion date, result, "
            "and recommended screening interval.\n\n"
            "Then simple date math determines urgency -- not AI guesswork. "
            "It takes the last completed date and the guideline timing window "
            "to determine when care is due, then measures how far past due that is today.\n\n"
            "The output: a prioritized, evidence-cited action list. What to do. Why. When. "
            "Where to route care. Every recommendation traced back to published clinical guidelines."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 18 — BestPath heals the Venn
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image15.png",
        notes_text=(
            "This holistically and fluidly addresses the systemic problem.\n\n"
            "The reality, in simplest terms, is that many people don't have a doctor, "
            "and doctors have too many patients.\n\n"
            "BestPath serves both of them."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 19 — Clinical Triage dashboard
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image5.png",
        notes_text=(
            "Two Modes, One Engine.\n\n"
            "Clinician mode -- Panel Management: A triage dashboard. "
            "Patients sorted into three columns:\n\n"
            "Red: Overdue + high risk -- act now.\n"
            "Yellow: Overdue + lower risk -- needs attention.\n"
            "Green: Up to date -- no action needed.\n\n"
            "Each split by confidence -- rich data vs. sparse. "
            "One click to approve an outreach with the patient's cited next step."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 20 — Care Navigator
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image10.png",
        notes_text=(
            "Patient mode -- Self-Service Navigator: For unattached Canadians. "
            "Enter your health information conversationally. Get your evidence-cited "
            "next step and routing to the right provider.\n\n"
            "Not just 'see a doctor.' A pharmacist can monitor blood pressure. "
            "A dietitian can help manage diabetes. A walk-in clinic can initiate a referral."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 21 — BestPath hammer (transition)
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image9.png",
        notes_text="[Pause -- let the image land.]")

    # ─────────────────────────────────────────────
    # SLIDE 22 — Why It Gets Adopted (text)
    # ─────────────────────────────────────────────
    add_text_slide(prs,
        title="Why It Gets Adopted",
        body_lines=[
            "**Invisible workflow:** No new systems to learn. No behavior change "
            "required. Patient data is automatically understood upon ingestion, "
            "and insights are automated.",
            "",
            "**Deterministic, auditable reasoning:** Clinicians see the guideline, "
            "the date math, the citation. Not a black box -- a transparent "
            "calculation they can trust and override.",
            "",
            "**Confidence flags:** When data is sparse or uncertain, BestPath "
            "says so. Low confidence, not guessed at.",
            "",
            "**Concise output:** One prioritized best action, not a wall of "
            "suggestions, delivered straight to the patient via text and email "
            "with one click.",
        ],
        notes_text=(
            "Why It Gets Adopted:\n\n"
            "Invisible workflow. No new systems to learn. No behavior change required. "
            "Patient data is automatically understood upon ingestion, and insights are automated.\n\n"
            "Deterministic, auditable reasoning. Clinicians see the guideline, the date math, "
            "the citation. Not a black box -- a transparent calculation they can trust and override.\n\n"
            "Confidence flags. When data is sparse or uncertain, BestPath says so. "
            "Low confidence, not guessed at.\n\n"
            "Concise output. One prioritized best action, not a wall of suggestions, "
            "delivered straight to the patient via text and email with one click."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 23 — Cute nail (BC alignment)
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image19.png",
        notes_text=(
            "BestPath is aligned with where BC is investing.\n\n"
            "BC is spending $2.8 billion on healthcare transformation over three years. "
            "BestPath fits directly into their stated priorities:\n\n"
            "Team-based primary care -- routes actions to the right provider role, "
            "not just physicians. Leverages pharmacists, nurses, and allied providers "
            "for preventive actions.\n\n"
            "Community and virtual care -- the patient navigator works without a clinic "
            "attachment, for all the Canadians without a primary care provider.\n\n"
            "Centralized digital systems -- architecture is FHIR-compatible, "
            "deployable against any EHR."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 24 — Hammer hitting nail (Impact)
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image6.png",
        notes_text=(
            "Impact.\n\n"
            "For providers: Saves clinical time by surfacing the highest-priority "
            "patients first. Prevents the 'who am I forgetting?' problem. "
            "Helps manage panels of over 1,000 patients.\n\n"
            "For patients: Improves access continuity. The system surfaces the "
            "specific relevant authoritative information in their best interest, "
            "and delivers it seamlessly. Unattached patients get a next step instead of nothing.\n\n"
            "For the system: Reduces avoidable ER reliance through better routing. "
            "Prioritizes overdue high-risk actions before they become emergencies, "
            "and routes people to the appropriate destination."
        ))

    # ─────────────────────────────────────────────
    # SLIDE 25 — FINAL: BestPath cityscape (DO NOT CHANGE)
    # ─────────────────────────────────────────────
    add_image_slide(prs, "image16.png",
        notes_text=(
            "BestPath makes proactive care operational:\n\n"
            "The right next step\n"
            "at the right time\n"
            "to the right place.\n\n"
            "It bridges patient data and public clinical evidence to deliver clear, "
            "cited actions that improve outcomes for people, patients, providers, "
            "and the health system.\n\n"
            "When we look at the data, we know who is at risk. The guidelines already "
            "say what should happen next. BestPath connects the two -- before it's too late.\n\n"
            "Thank you."
        ))

    out = os.path.join(BASE, "BestPath-Presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
